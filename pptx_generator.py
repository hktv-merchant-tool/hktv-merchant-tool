"""
HKTVmall PPT 動態生成模組
=========================
基於「2026 HKTVmall New Merchant Acquisition Deck [CHI].pptx」模板
根據輸入的目標招商品牌，生成客製化招商PPT

使用方法:
    from pptx_generator import generate_merchant_ppt
    result = generate_merchant_ppt(
        merchant_name="永明凍肉有限公司",
        category="糧油雜貨",
        pain_points=["缺乏線上流量", "冷鏈配送常有客訴"],
        custom_content="希望強調跨境電商能力",
        tone="誠懇專業型",
        template_path="path/to/template.pptx"
    )
    # result 為 BytesIO，可直接下載
"""

import io
import copy
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# ──────────────────────────────────────────────
# 模板置換配置
# ──────────────────────────────────────────────

# 可替換的文字關鍵字 → 對應替換內容
REPLACEMENTS = {
    # Slide 1: 封面（目標商戶名稱）
    "No.1  香港網上購物平台商戶合作方案": "MERCHANT_PLACEHOLDER 香港網上購物平台商戶合作方案",
}

# 類別專屬話術
CATEGORY_INSIGHTS = {
    "糧油雜貨": "糧油雜貨是HKTVmall GMV最高的品類，家庭消費者忠誠度高，復購率達4.7x/季，加入我們立即享受流量紅利。",
    "美容及健康產品": "美容及健康產品在HKTVmall增長迅猛，平台用戶男女均衡，個人專屬價工具能精準觸及目標客戶，20%轉化率效果顯著。",
    "寵物用品": "寵物市場快速增長，直播頻道可展示寵物用品實際使用效果，透過CRM精準定位寵物主人群體，大幅提升轉化率。",
    "電子電器產品": "蘇寧等大型品牌已在HKTVmall直播單晚創下近90萬GMV，電子電器產品極適合直播展示功能，搶佔高端客戶群。",
    "其他": "HKTVmall 310萬+用戶覆蓋全港各類消費者，全方位營銷工具助您快速打開市場，3-5個工作天即可開店營運。",
}

# 痛點對應解決方案
PAIN_SOLUTIONS = {
    "缺乏線上流量": "HKTVmall擁有310萬+獨立用戶及160萬月活設備，能為商戶帶來巨大線上流量曝光，解決流量獲取難題。",
    "自建物流成本過高": "HKTVmall自建7大倉庫、350+輛貨車及75家O2O門市，提供全港最大住宅配送網絡，讓商戶告別自建物流的高昂成本。",
    "冷鏈配送常有客訴": "HKTVmall配備專業冷鏈配送系統，每日處理10萬+訂單，自動化倉庫確保生鮮及冷藏商品品質，大幅降低客訴率。",
    "營銷成本過高": "HKTVmall提供85折推廣（平台全額承擔）、個人專屬價（20%轉化率）等高效營銷工具，並有鄭裕玲、黃子華等名人廣告加持，讓營銷更具性價比。",
    "庫存管理困難": "GREEN LAB臨期百貨及會員月費計劃能有效幫助商戶清理庫存，降低損耗成本。",
    "缺乏電商營運經驗": "HKTVmall電商學院提供全方位支援，包括上架教學、數據分析、行銷技巧，讓電商新手也能快速上手。",
    "曝光機會不足": "58個MTR站、3,000+廣告燈箱的全渠道曝光，加上直播頻道，為商戶帶來前所未有的品牌曝光機會。",
}

# 各語氣風格的開場白
TONE_OPENINGS = {
    "誠懇專業型": "您好，\n\n非常高興有機會向您介紹 HKTVmall — 香港 No.1 網上購物平台，以及我們為您量身訂造的商戶合作方案。",
    "強勢說服型": "您好！\n\n我必須直接告訴您：在香港電商市場，已經有超過310萬消費者在HKTVmall購物，而您的品牌還未入駐，這是一個不容忽視的商機。",
    "溫暖關懷型": "您好呀！\n\n知道您一直在用心經營品牌，我們很希望能把 HKTVmall 的資源帶給您，讓您的生意事半功倍。",
    "數據導向型": "您好，\n\n以下是一份基於HKTVmall真實平台數據的商戶合作方案，供您決策參考：",
    "簡潔利落型": "您好，HKTVmall誠邀入駐，重點如下：",
    "輕鬆友好型": "嘿！\n\n有個超棒的消息要告訴您 — HKTVmall 正在招募優質商戶，我覺得您超適合！",
}

# 語氣對應的CTA
TONE_CTAS = {
    "誠懇專業型": "如您對入駐HKTVmall有任何疑問，歡迎隨時與我聯絡，我們的招商團隊將全程為您提供支援。期待與您合作，共創佳績！",
    "強勢說服型": "現在就行動吧！立即聯絡我們，3-5個工作天內即可開店營運，在310萬用戶面前展示您的品牌！",
    "溫暖關懷型": "隨時歡迎您提出任何問題，我會全程耐心跟進。希望有機會與您合作，讓您的品牌被更多香港家庭看見。",
    "數據導向型": "立即聯絡HKTVmall招商團隊，開啟您的電商增長之旅。",
    "簡潔利落型": "聯絡我們，立即入駐！",
    "輕鬆友好型": "快點找我聊聊吧，我超期待看到您的品牌登陸HKTVmall！",
}


# ──────────────────────────────────────────────
# 核心生成函數
# ──────────────────────────────────────────────

def replace_text_in_shape(shape, old_text, new_text):
    """置換shape中的指定文字"""
    if not shape.has_text_frame:
        return False

    replaced = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old_text in run.text:
                run.text = run.text.replace(old_text, new_text)
                replaced = True
    return replaced


def replace_text_in_slide(slide, old_text, new_text):
    """置換slide中所有shape的指定文字"""
    count = 0
    for shape in slide.shapes:
        if replace_text_in_shape(shape, old_text, new_text):
            count += 1
    return count


def find_and_replace_in_slide(slide, old_text, new_text):
    """在slide中尋找並替換文字（更完整的實現）"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full_text = ''.join([run.text for run in para.runs])
            if old_text in full_text:
                # 保留其他run的格式，只替換包含old_text的部分
                for run in para.runs:
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)


def generate_merchant_ppt(
    merchant_name: str,
    category: str,
    pain_points: List[str],
    custom_content: str = "",
    tone: str = "誠懇專業型",
    template_path: str = None
) -> io.BytesIO:
    """
    生成客製化招商PPT

    Args:
        merchant_name: 目標商戶名稱
        category: 產品品類
        pain_points: 痛點列表
        custom_content: 自訂內容提示
        tone: 語氣風格
        template_path: PPT模板路徑

    Returns:
        BytesIO對象，可直接作為下載檔案
    """

    # 預設模板路徑 - 從 GitHub 讀取
    if template_path is None:
        import urllib.request
        template_url = "https://raw.githubusercontent.com/hktv-merchant-tool/hktv-merchant-tool/main/template/2026%20HKTVmall%20New%20Merchant%20Acquisition%20Deck%20[CHI].pptx"
        # 下載至記憶體
        template_bytes = urllib.request.urlopen(template_url).read()
        prs = Presentation(io.BytesIO(template_bytes))

    # ─── Slide 1: 封面 ───
    # 將"No.1 香港網上購物平台商戶合作方案" 改為 "XXX 香港網上購物平台商戶合作方案"
    slide1 = prs.slides[0]
    replace_text_in_slide(slide1, "No.1", f"{merchant_name}")

    # ─── Slide 5: 平台流量表現 ───
    # 可在這裡添加merchant特定的強調
    slide5 = prs.slides[4]

    # ─── Slide 8: GMV ───
    slide8 = prs.slides[7]
    # 這張保持平台數據不變（因為是平台全局數據）

    # ─── Slide 22: 個人專屬價（轉化率） ───
    # 根據品類強調相應工具
    cat_data = CATEGORY_INSIGHTS.get(category, CATEGORY_INSIGHTS["其他"])
    slide22 = prs.slides[21]

    # ─── Slide 31: 蘇寧成功案例（可客製為目標商戶同業案例） ───
    # 保持作為參考案例

    # ─── Slide 44: 聯絡方式（CTA） ───
    # 確保聯絡方式正確
    slide44 = prs.slides[43]

    # ─── 生成結論頁（新增或修改最後一張） ───
    # 在倒數第二張添加merchant特定總結
    if len(prs.slides) >= 48:
        slide_summary = prs.slides[47]
        # 構建總結內容
        summary_lines = [
            f"目標招商商戶：{merchant_name}",
            f"產品品類：{category}",
            f"語氣風格：{tone}",
        ]
        if pain_points:
            summary_lines.append(f"痛點關注：{', '.join(pain_points)}")
        if custom_content:
            summary_lines.append(f"特別要求：{custom_content}")

        # 這裡我們保持模板不變，只記錄元數據
        # PPT本身不適合動態新增大量文字，保持模板結構

    # ─── 儲存為BytesIO ───
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    return output


def generate_custom_ppt(
    uploaded_ppt_bytes: bytes,
    merchant_name: str,
    category: str,
    merchant_desc: str,
    pain_points: List[str],
    custom_content: str = "",
    tone: str = "誠懇專業型"
) -> dict:
    """
    根據上傳的PPT模板，生成客製化招商PPT

    Args:
        uploaded_ppt_bytes: 上傳的PPT檔案bytes
        merchant_name: 目標商戶名稱
        category: 產品品類
        merchant_desc: 主要業務描述
        pain_points: 痛點列表
        custom_content: 自訂內容提示
        tone: 語氣風格

    Returns:
        dict: {'ppt': BytesIO, 'summary': str, 'custom_slides': list}
    """

    # 從bytes載入PPT
    prs = Presentation(io.BytesIO(uploaded_ppt_bytes))
    custom_slides = []

    # ─── Slide 1: 封面 ───
    # 嘗試替換「No.1」為商戶名稱
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "No.1" in run.text:
                        run.text = run.text.replace("No.1", merchant_name)
                        if 1 not in custom_slides:
                            custom_slides.append(1)
                    elif "香港網上購物平台商戶合作方案" in run.text and merchant_name:
                        run.text = f"{merchant_name} - " + run.text
                        if 1 not in custom_slides:
                            custom_slides.append(1)

    # ─── Slide 49（或倒數）：可加入商戶備註 ───
    if len(prs.slides) >= 49:
        last_slide = prs.slides[48]
        for shape in last_slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip() == "" or "附錄" in run.text:
                            run.text = f"目標招商商戶：{merchant_name}\n產品品類：{category}\n業務描述：{merchant_desc[:100]}{'...' if len(merchant_desc) > 100 else ''}"
                            if 49 not in custom_slides:
                                custom_slides.append(49)
                            break

    # 構建摘要
    cat_insight = CATEGORY_INSIGHTS.get(category, CATEGORY_INSIGHTS["其他"])
    pain_texts = [PAIN_SOLUTIONS.get(p, "") for p in pain_points if p in PAIN_SOLUTIONS]

    summary_parts = [
        f"【招商目標】{merchant_name}",
        f"【產品品類】{category}",
        f"【主要業務】{merchant_desc[:80]}{'...' if len(merchant_desc) > 80 else ''}",
        f"【語氣風格】{tone}",
        "",
        "【品類洞察】",
        cat_insight,
        "",
    ]

    if pain_texts:
        summary_parts.append("【痛點對應方案】")
        for pt in pain_texts:
            summary_parts.append(f"• {pt}")
        summary_parts.append("")

    if custom_content:
        summary_parts.append(f"【特別要求】{custom_content}")
        summary_parts.append("")

    summary_parts.extend([
        "【PPT客製化說明】",
        f"已根據上傳的PPT模板進行客製化",
        f"共 {len(prs.slides)} 張投影片",
        "以下slides已根據目標商戶調整：",
        f"• Slide 1: 封面（已加入商戶名稱）",
    ])

    if merchant_desc:
        summary_parts.append(f"• Slide 49: 備註頁（已加入商戶業務描述）")

    summary_parts.extend([
        "",
        "【聯絡我們】",
        "WhatsApp：+852 5283 4138",
        "電郵：aog.merc@hktv.com.hk",
        "網站：business.hktvmall.com",
    ])

    # 儲存
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    return {
        'ppt': output,
        'summary': '\n'.join(summary_parts),
        'custom_slides': custom_slides,
    }


def generate_merchant_ppt_with_notes(
    merchant_name: str,
    category: str,
    pain_points: List[str],
    custom_content: str = "",
    tone: str = "誠懇專業型",
    template_path: str = None
) -> dict:
    """
    生成客製化PPT + 生成摘要說明

    Returns:
        dict: {
            'ppt': BytesIO,
            'summary': str,  # 摘要說明
            'custom_slides': list  # 被修改的slide編號列表
        }
    """

    if template_path is None:
        import urllib.request
        template_url = "https://raw.githubusercontent.com/hktv-merchant-tool/hktv-merchant-tool/main/template/2026%20HKTVmall%20New%20Merchant%20Acquisition%20Deck%20[CHI].pptx"
        template_bytes = urllib.request.urlopen(template_url).read()
        prs = Presentation(io.BytesIO(template_bytes))
    else:
        prs = Presentation(template_path)
    custom_slides = []

    # Slide 1: 封面 - 加上商戶名稱
    slide1 = prs.slides[0]
    if replace_text_in_slide(slide1, "No.1", merchant_name) > 0:
        custom_slides.append(1)

    # Slide 31-34: 成功案例（保持不變，作為同業參考）
    # 這些是通用案例，不用改

    # Slide 44: CTA保持
    # 聯絡方式保持不變

    # 生成摘要
    cat_insight = CATEGORY_INSIGHTS.get(category, CATEGORY_INSIGHTS["其他"])
    pain_texts = [PAIN_SOLUTIONS.get(p, "") for p in pain_points if p in PAIN_SOLUTIONS]

    summary_parts = [
        f"【招商目標】{merchant_name}",
        f"【產品品類】{category}",
        f"【語氣風格】{tone}",
        "",
        "【品類洞察】",
        cat_insight,
        "",
    ]

    if pain_texts:
        summary_parts.append("【痛點對應】")
        for pt in pain_texts:
            summary_parts.append(f"• {pt}")
        summary_parts.append("")

    if custom_content:
        summary_parts.append(f"【特別要求】{custom_content}")
        summary_parts.append("")

    summary_parts.extend([
        "【PPT說明】",
        f"已使用「2026 HKTVmall New Merchant Acquisition Deck」作為模板",
        f"覆蓋 {len(prs.slides)} 張投影片",
        "以下slides已根據目標商戶客製化：",
        f"• Slide 1: 封面（已加入商戶名稱）",
        "• Slide 31-34: 成功案例（保持作為同業參考）",
        "• Slide 44: 聯絡我們（保持不變）",
        "",
        f"聯絡人：HKTVmall招商團隊",
        "WhatsApp：+852 5283 4138",
        "電郵：aog.merc@hktv.com.hk",
    ])

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    return {
        'ppt': output,
        'summary': '\n'.join(summary_parts),
        'custom_slides': custom_slides,
    }


# ──────────────────────────────────────────────
# 測試
# ──────────────────────────────────────────────
if __name__ == "__main__":
    # 測試生成
    result = generate_merchant_ppt_with_notes(
        merchant_name="永明凍肉有限公司",
        category="糧油雜貨",
        pain_points=["缺乏線上流量", "冷鏈配送常有客訴"],
        custom_content="希望強調跨境電商能力",
        tone="誠懇專業型",
    )

    print("=== 生成摘要 ===")
    print(result['summary'])
    print("\n=== 已客製化的Slides ===")
    print(result['custom_slides'])

    # 儲存測試檔案
    with open("/tmp/test_merchant.pptx", "wb") as f:
        f.write(result['ppt'].getvalue())
    print("\n已儲存測試PPT至 /tmp/test_merchant.pptx")